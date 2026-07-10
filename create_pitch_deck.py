import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Dark Theme / Stellar Brand Colors)
BG_COLOR = RGBColor(8, 11, 17)        # Very Dark Slate/Black
CARD_COLOR = RGBColor(16, 22, 35)     # Glassmorphic Slate/Navy
BORDER_COLOR = RGBColor(30, 41, 59)   # Cool Slate Gray
TEXT_WHITE = RGBColor(243, 244, 246)  # Gray 100
TEXT_MUTED = RGBColor(156, 163, 175)  # Gray 400
ACCENT_BLUE = RGBColor(59, 130, 246)  # Stellar/Soroban Blue
ACCENT_GOLD = RGBColor(245, 158, 11)  # Lumen Amber/Gold
ACCENT_PURPLE = RGBColor(139, 92, 246) # Client Purple

blank_slide_layout = prs.slide_layouts[6]

def apply_background(slide):
    # Create background fill
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background() # No border
    
    # Accent top border line
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.1))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_BLUE
    accent_bar.line.fill.background()

def add_header(slide, title, category="STELLAR ESCROW PITCH DECK"):
    # Category Tracker
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.3))
    cat_tf = cat_box.text_frame
    cat_tf.word_wrap = True
    cat_p = cat_tf.paragraphs[0]
    cat_p.text = category.upper()
    cat_p.font.name = "Calibri"
    cat_p.font.size = Pt(10)
    cat_p.font.bold = True
    cat_p.font.color.rgb = ACCENT_GOLD
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.name = "Trebuchet MS"
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = TEXT_WHITE

def draw_card(slide, left, top, width, height, title, body_paragraphs, accent_color=None):
    # Card Background shape
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_COLOR
    card.line.color.rgb = accent_color if accent_color else BORDER_COLOR
    card.line.width = Pt(1.5)
    
    # Card Content Textbox
    textbox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = textbox.text_frame
    tf.word_wrap = True
    
    # Card Title
    title_p = tf.paragraphs[0]
    title_p.text = title
    title_p.font.name = "Trebuchet MS"
    title_p.font.size = Pt(18)
    title_p.font.bold = True
    title_p.font.color.rgb = accent_color if accent_color else ACCENT_BLUE
    
    # Card Body
    for text in body_paragraphs:
        p = tf.add_paragraph()
        p.text = text
        p.font.name = "Calibri"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(8)

# ----------------- SLIDE 1: Title Slide -----------------
slide1 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide1)

# Large Center Decorative Shape
bg_glow = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.66), Inches(1.5), Inches(6), Inches(4.5))
bg_glow.fill.solid()
bg_glow.fill.fore_color.rgb = RGBColor(15, 23, 42)
bg_glow.line.fill.background()

# Title text block
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

p_main = tf.paragraphs[0]
p_main.alignment = PP_ALIGN.CENTER
p_main.text = "StellarEscrow"
p_main.font.name = "Trebuchet MS"
p_main.font.size = Pt(64)
p_main.font.bold = True
p_main.font.color.rgb = TEXT_WHITE

p_sub = tf.add_paragraph()
p_sub.alignment = PP_ALIGN.CENTER
p_sub.text = "Trustless milestone-based payments for the global freelance economy."
p_sub.font.name = "Calibri"
p_sub.font.size = Pt(22)
p_sub.font.color.rgb = ACCENT_GOLD
p_sub.space_before = Pt(14)

p_built = tf.add_paragraph()
p_built.alignment = PP_ALIGN.CENTER
p_built.text = "BUILT ON STELLAR  •  POWERED BY SOROBAN"
p_built.font.name = "Calibri"
p_built.font.size = Pt(12)
p_built.font.bold = True
p_built.font.color.rgb = ACCENT_BLUE
p_built.space_before = Pt(20)


# ----------------- SLIDE 2: The Problem -----------------
slide2 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide2)
add_header(slide2, "The Freelance Payment & Trust Gap")

# 4 Stat Cards
draw_card(slide2, Inches(0.8), Inches(2.0), Inches(2.6), Inches(4.3), 
          "Payment Delays", 
          ["29% of freelancers experience chronic payment delays.", "Clients hold leverage, leaving freelancers with zero protection after work is submitted."], 
          ACCENT_GOLD)

draw_card(slide2, Inches(3.7), Inches(2.0), Inches(2.6), Inches(4.3), 
          "Heavy Payout Fees", 
          ["Cross-border payments cost 5% to 10% in currency conversion and wire transfer fees.", "Traditional banks and payment networks exploit remote builders."],
          ACCENT_GOLD)

draw_card(slide2, Inches(6.6), Inches(2.0), Inches(2.6), Inches(4.3), 
          "Clearing Delays", 
          ["International wire clearings take 3 to 7 business days.", "Slow liquidity cycles restrict cash flows for digital creators."],
          ACCENT_GOLD)

draw_card(slide2, Inches(9.5), Inches(2.0), Inches(2.6), Inches(4.3), 
          "Marketplace Cut", 
          ["Upwork & Fiverr extract up to 20% in platform commissions.", "Middlemen capture the value created entirely by independent talent."],
          ACCENT_GOLD)


# ----------------- SLIDE 3: The Solution -----------------
slide3 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide3)
add_header(slide3, "Trustless Milestone-Based Escrows")

# Left Column (Core features)
left_box = slide3.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.0), Inches(4.5))
left_tf = left_box.text_frame
left_tf.word_wrap = True

features = [
    ("Milestone-Locked Funds", "Clients secure payments upfront directly inside an on-chain smart contract escrow. Freelancers verify the capital is locked before starting work."),
    ("Automatic Payout Releases", "Upon task completion and client approval, the escrowed funds instantly unlock on-chain, eliminating middleman dependency."),
    ("Zero Brokerage Commissions", "No platform commissions or escrow service cut. Value transfers directly from Client to Freelancer, trustlessly.")
]

for i, (title, desc) in enumerate(features):
    p_t = left_tf.add_paragraph() if i > 0 else left_tf.paragraphs[0]
    p_t.text = f"✔ {title}"
    p_t.font.name = "Trebuchet MS"
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = ACCENT_BLUE
    p_t.space_before = Pt(15) if i > 0 else Pt(0)
    
    p_d = left_tf.add_paragraph()
    p_d.text = desc
    p_d.font.name = "Calibri"
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = TEXT_WHITE
    p_d.space_before = Pt(4)

# Right Column (Visual flow diagram boxes)
# Card 1: Client
draw_card(slide3, Inches(6.5), Inches(2.0), Inches(5.5), Inches(1.2), 
          "1. Client Deposits XLM", 
          ["Locks project milestone amount into the Soroban contract."], ACCENT_PURPLE)

# Down Arrow
arrow = slide3.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.0), Inches(3.3), Inches(0.5), Inches(0.4))
arrow.fill.solid()
arrow.fill.fore_color.rgb = BORDER_COLOR
arrow.line.fill.background()

# Card 2: Contract
draw_card(slide3, Inches(6.5), Inches(3.8), Inches(5.5), Inches(1.2), 
          "2. Secure Soroban Escrow", 
          ["Contract holds funds securely on-chain. Work is submitted."], ACCENT_BLUE)

# Down Arrow
arrow2 = slide3.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.0), Inches(5.1), Inches(0.5), Inches(0.4))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = BORDER_COLOR
arrow2.line.fill.background()

# Card 3: Freelancer
draw_card(slide3, Inches(6.5), Inches(5.6), Inches(5.5), Inches(1.2), 
          "3. Automated Release", 
          ["Client approves the milestone, triggering auto-payment transfer."], ACCENT_GOLD)


# ----------------- SLIDE 4: Why Stellar -----------------
slide4 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide4)
add_header(slide4, "Why Build on Stellar & Soroban?")

draw_card(slide4, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5), 
          "Extreme Efficiency", 
          ["• 3-5 Second Payouts: Settles transactions near-instantly compared to days via traditional banking networks.", 
           "• Fractional Fees: Transaction costs are less than 1/10,000th of a cent, avoiding Ethereum's gas fee surges.",
           "• High Scalability: Built to handle global retail payment throughput natively."], 
          ACCENT_BLUE)

draw_card(slide4, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.5), 
          "Ecosystem Anchors", 
          ["• Direct Fiat Access: Stellar Anchors allow direct USD/Euro stablecoin to local bank conversions.", 
           "• Global Reach: Enables localized remittance and withdrawal rails across 180+ countries.",
           "• Embedded Compliance: Supports regulated assets and identity hooks natively on-chain."], 
          ACCENT_BLUE)

draw_card(slide4, Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.5), 
          "Soroban Smart Contracts", 
          ["• Secure Execution: Rust-based WebAssembly (Wasm) architecture provides sandboxed safety.", 
           "• Minimal Fees: State archive storage costs are optimized, keeping smart contract state fees low.",
           "• Modular Escrows: Highly extensible milestone dispute resolution logic can be implemented."], 
          ACCENT_BLUE)


# ----------------- SLIDE 5: Product Demo -----------------
slide5 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide5)
add_header(slide5, "Intuitive Web3 User Experience")

# Image 1: Dashboard UI
try:
    slide5.shapes.add_picture("image-6.png", Inches(0.8), Inches(2.0), Inches(3.6), Inches(2.3))
except Exception as e:
    draw_card(slide5, Inches(0.8), Inches(2.0), Inches(3.6), Inches(2.3), "Dashboard UI Image", [str(e)], BORDER_COLOR)

draw_card(slide5, Inches(0.8), Inches(4.5), Inches(3.6), Inches(2.0), 
          "1. Dashboard Overviews", 
          ["The client/freelancer sees active escrow milestones, on-chain balances, and custom visual statuses."], 
          BORDER_COLOR)

# Image 2: Onboarding Flow
try:
    slide5.shapes.add_picture("image-5.png", Inches(4.8), Inches(2.0), Inches(3.6), Inches(2.3))
except Exception as e:
    draw_card(slide5, Inches(4.8), Inches(2.0), Inches(3.6), Inches(2.3), "Onboarding Flow Image", [str(e)], BORDER_COLOR)

draw_card(slide5, Inches(4.8), Inches(4.5), Inches(3.6), Inches(2.0), 
          "2. Guided Onboarding", 
          ["Interactive tutorial directs first-time Web3 users to install Freighter and get test XLM from Friendbot."], 
          BORDER_COLOR)

# Image 3: PostHog / Telemetry
try:
    slide5.shapes.add_picture("image-7.png", Inches(8.8), Inches(2.0), Inches(3.6), Inches(2.3))
except Exception as e:
    draw_card(slide5, Inches(8.8), Inches(2.0), Inches(3.6), Inches(2.3), "Telemetry Image", [str(e)], BORDER_COLOR)

draw_card(slide5, Inches(8.8), Inches(4.5), Inches(3.6), Inches(2.0), 
          "3. Verified Telemetry", 
          ["Integrated event trackers monitor active user funnels and Soroban ledger transactions live."], 
          BORDER_COLOR)


# ----------------- SLIDE 6: Architecture -----------------
slide6 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide6)
add_header(slide6, "Decentralized Architecture Overview")

# Architecture Flow chart
# Box 1: React UI
draw_card(slide6, Inches(0.8), Inches(2.2), Inches(2.5), Inches(3.5), 
          "1. React Frontend", 
          ["• Tailwind CSS for UI layout", "• Glassmorphic design system", "• Context API state management"], 
          ACCENT_PURPLE)

# Right Arrow
r_arrow1 = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.4), Inches(3.7), Inches(0.4), Inches(0.4))
r_arrow1.fill.solid()
r_arrow1.fill.fore_color.rgb = BORDER_COLOR
r_arrow1.line.fill.background()

# Box 2: Wallet Kit
draw_card(slide6, Inches(3.9), Inches(2.2), Inches(2.5), Inches(3.5), 
          "2. Wallet Connect", 
          ["• StellarWalletsKit SDK", "• Freighter browser extension", "• Secure transaction signing"], 
          ACCENT_BLUE)

# Right Arrow
r_arrow2 = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(3.7), Inches(0.4), Inches(0.4))
r_arrow2.fill.solid()
r_arrow2.fill.fore_color.rgb = BORDER_COLOR
r_arrow2.line.fill.background()

# Box 3: Smart Contract
draw_card(slide6, Inches(7.0), Inches(2.2), Inches(2.5), Inches(3.5), 
          "3. Soroban Escrow", 
          ["• Rust Smart Contract", "• Milestone state mapping", "• State-locked escrow balance"], 
          ACCENT_BLUE)

# Right Arrow
r_arrow3 = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.6), Inches(3.7), Inches(0.4), Inches(0.4))
r_arrow3.fill.solid()
r_arrow3.fill.fore_color.rgb = BORDER_COLOR
r_arrow3.line.fill.background()

# Box 4: Ledger
draw_card(slide6, Inches(10.1), Inches(2.2), Inches(2.5), Inches(3.5), 
          "4. Stellar Ledger", 
          ["• Decentralized ledger", "• Fast consensus", "• Public audit explorer"], 
          ACCENT_GOLD)

# Telemetry note below
tel_box = slide6.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11.8), Inches(0.8))
tel_tf = tel_box.text_frame
tel_tf.word_wrap = True
tel_p = tel_tf.paragraphs[0]
tel_p.text = "⚡ STABILITY & TELEMETRY LAYER: Real-time user event indexing via PostHog. Frontend exceptions captured and logged in real-time by Sentry to ensure 99.9% client-side uptime."
tel_p.font.name = "Calibri"
tel_p.font.size = Pt(11)
tel_p.font.bold = True
tel_p.font.color.rgb = TEXT_MUTED


# ----------------- SLIDE 7: Market Opportunity -----------------
slide7 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide7)
add_header(slide7, "Market Opportunity: Global Freelance Economy")

draw_card(slide7, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), 
          "A Massively Expanding Global Market", 
          [
              "• Market Size: Projections reach $455 Billion by Mastercard gig economy studies.",
              "• Remote Work Adoption: The transition to remote-first collaboration has amplified cross-border hiring.",
              "• Transaction Overhead: Over $100 Billion is spent annually in bank fees and marketplace cuts.",
              "• Maturing Rails: Web3 wallet usage and native currency adoption are hitting mainstream thresholds."
          ], 
          ACCENT_BLUE)

draw_card(slide7, Inches(6.8), Inches(2.0), Inches(5.5), Inches(4.5), 
          "The Cross-Border Payment Bottleneck", 
          [
              "• High Friction: Global transfers face currency conversion markups, regulatory delays, and intermediary banking cuts.",
              "• Trust Bottleneck: Escrows are either highly expensive (marketplaces extracting 20%) or non-existent.",
              "• Opportunity: StellarEscrow offers trustless security for the same price as fractions of a cent, removing intermediary margin extraction entirely."
          ], 
          ACCENT_GOLD)


# ----------------- SLIDE 8: Traction So Far -----------------
slide8 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide8)
add_header(slide8, "Testnet Validation & Metrics")

# Large Stat Callouts
# Stat 1
s1 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(3.6), Inches(2.0))
s1.fill.solid()
s1.fill.fore_color.rgb = CARD_COLOR
s1.line.color.rgb = BORDER_COLOR
s1_tf = s1.text_frame
s1_tf.word_wrap = True
s1_p1 = s1_tf.paragraphs[0]
s1_p1.text = "55+"
s1_p1.font.name = "Trebuchet MS"
s1_p1.font.size = Pt(44)
s1_p1.font.bold = True
s1_p1.font.color.rgb = ACCENT_GOLD
s1_p1.alignment = PP_ALIGN.CENTER
s1_p2 = s1_tf.add_paragraph()
s1_p2.text = "Onboarded Testnet Users"
s1_p2.font.name = "Calibri"
s1_p2.font.size = Pt(14)
s1_p2.font.color.rgb = TEXT_WHITE
s1_p2.alignment = PP_ALIGN.CENTER

# Stat 2
s2 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(2.2), Inches(3.6), Inches(2.0))
s2.fill.solid()
s2.fill.fore_color.rgb = CARD_COLOR
s2.line.color.rgb = BORDER_COLOR
s2_tf = s2.text_frame
s2_tf.word_wrap = True
s2_p1 = s2_tf.paragraphs[0]
s2_p1.text = "72+"
s2_p1.font.name = "Trebuchet MS"
s2_p1.font.size = Pt(44)
s2_p1.font.bold = True
s2_p1.font.color.rgb = ACCENT_BLUE
s2_p1.alignment = PP_ALIGN.CENTER
s2_p2 = s2_tf.add_paragraph()
s2_p2.text = "On-Chain Transactions Logged"
s2_p2.font.name = "Calibri"
s2_p2.font.size = Pt(14)
s2_p2.font.color.rgb = TEXT_WHITE
s2_p2.alignment = PP_ALIGN.CENTER

# Stat 3
s3 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(2.2), Inches(3.6), Inches(2.0))
s3.fill.solid()
s3.fill.fore_color.rgb = CARD_COLOR
s3.line.color.rgb = BORDER_COLOR
s3_tf = s3.text_frame
s3_tf.word_wrap = True
s3_p1 = s3_tf.paragraphs[0]
s3_p1.text = "4.8 / 5"
s3_p1.font.name = "Trebuchet MS"
s3_p1.font.size = Pt(44)
s3_p1.font.bold = True
s3_p1.font.color.rgb = ACCENT_PURPLE
s3_p1.alignment = PP_ALIGN.CENTER
s3_p2 = s3_tf.add_paragraph()
s3_p2.text = "Average User Satisfaction Score"
s3_p2.font.name = "Calibri"
s3_p2.font.size = Pt(14)
s3_p2.font.color.rgb = TEXT_WHITE
s3_p2.alignment = PP_ALIGN.CENTER

# Key Feedback Callout Box
feedback_card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.7), Inches(11.6), Inches(1.8))
feedback_card.fill.solid()
feedback_card.fill.fore_color.rgb = CARD_COLOR
feedback_card.line.color.rgb = ACCENT_BLUE
feedback_tf = feedback_card.text_frame
feedback_tf.word_wrap = True

f_title = feedback_tf.paragraphs[0]
f_title.text = "🗣️ KEY USER INSIGHT"
f_title.font.name = "Trebuchet MS"
f_title.font.size = Pt(16)
f_title.font.bold = True
f_title.font.color.rgb = ACCENT_BLUE

f_desc = feedback_tf.add_paragraph()
f_desc.text = '"Having milestone funds locked transparently on-chain gives freelancers immediate peace of mind before writing a line of code, while clients love only paying for verified, complete deliverables."'
f_desc.font.name = "Calibri"
f_desc.font.size = Pt(13)
f_desc.font.color.rgb = TEXT_WHITE
f_desc.font.italic = True
f_desc.space_before = Pt(8)


# ----------------- SLIDE 9: Growth Strategy -----------------
slide9 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide9)
add_header(slide9, "Go-To-Market & Growth Strategy")

draw_card(slide9, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.1), 
          "Freelance Communities Integration", 
          ["Targeting freelance subreddits, Discord channels, and LinkedIn contract networks.", "Grassroots marketing focused on non-payment horror stories and heavy marketplace cuts."], 
          BORDER_COLOR)

draw_card(slide9, Inches(6.8), Inches(2.0), Inches(5.5), Inches(2.1), 
          "Pilot Agency Programs", 
          ["Onboarding boutique software development and design agencies as power users.", "Focus on agency clients who value structured milestone billing."], 
          BORDER_COLOR)

draw_card(slide9, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.1), 
          "Built-in Viral Invite Loop", 
          ["When a client creates a project, the freelancer receives an email invite link to claim it.", "Dual-sided onboarding: each user introduces another to the platform."], 
          ACCENT_BLUE)

draw_card(slide9, Inches(6.8), Inches(4.5), Inches(5.5), Inches(2.1), 
          "Educational Content Funnel", 
          ["Writing clear, jargon-free guides on Web3 utility for traditional remote workers.", "Hiding blockchain complexity beneath clean, familiar interfaces."], 
          BORDER_COLOR)


# ----------------- SLIDE 10: Roadmap -----------------
slide10 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide10)
add_header(slide10, "Product Roadmap & Future Milestones")

# Horizontal Roadmap
draw_card(slide10, Inches(0.8), Inches(2.5), Inches(2.6), Inches(3.8), 
          "Phase 1: Testnet MVP", 
          ["• Complete milestone escrow contract", "• StellarWalletsKit interface", "• Onboarded 50+ real testers", "• Sentry + PostHog integrated"], 
          ACCENT_BLUE)

draw_card(slide10, Inches(3.7), Inches(2.5), Inches(2.6), Inches(3.8), 
          "Phase 2: Anchor Connect", 
          ["• Stellar Anchor integrations", "• Local banking on/off-ramps", "• Support for fiat stablecoins", "• Regulated identity profiles"], 
          BORDER_COLOR)

draw_card(slide10, Inches(6.6), Inches(2.5), Inches(2.6), Inches(3.8), 
          "Phase 3: Escrow DAO", 
          ["• Multi-arbitrator dispute resolution", "• Decentralized juror pools", "• Subscription payment options", "• Retainer lock contracts"], 
          BORDER_COLOR)

draw_card(slide10, Inches(9.5), Inches(2.5), Inches(2.6), Inches(3.8), 
          "Phase 4: Mainnet Payouts", 
          ["• Mainnet contract deploy", "• Viral growth loop incentives", "• Integration with global invoice apps", "• Scale to $10M+ transaction volume"], 
          ACCENT_GOLD)


# ----------------- SLIDE 11: Team / Ask -----------------
slide11 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide11)
add_header(slide11, "Ecosystem Integration & Ask")

draw_card(slide11, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), 
          "The Ask: Growth & Mainnet Support", 
          [
              "• Ecosystem Funding: Development grant to subsidize initial mainnet fee reserves.",
              "• Design Partners: Boutique agencies to test fiat on-ramp integrations.",
              "• Feedback & Reviews: Join our testnet trials to refine smart contract escrow usability."
          ], 
          ACCENT_GOLD)

draw_card(slide11, Inches(6.8), Inches(2.0), Inches(5.5), Inches(4.5), 
          "Ecosystem Integration Opportunities", 
          [
              "• Invoice Software API: Plug trustless escrows directly into standard billing systems.",
              "• Freelance Marketplace Plugins: Add 'Pay via StellarEscrow' buttons to external platforms.",
              "• Stellar Network Alignment: Drive transaction volume and anchor utility on Stellar Testnet and Mainnet."
          ], 
          BORDER_COLOR)


# ----------------- SLIDE 12: Closing / Contact Slide -----------------
slide12 = prs.slides.add_slide(blank_slide_layout)
apply_background(slide12)

# Large Center Decorative Shape
bg_glow2 = slide12.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.66), Inches(1.5), Inches(6), Inches(4.5))
bg_glow2.fill.solid()
bg_glow2.fill.fore_color.rgb = RGBColor(15, 23, 42)
bg_glow2.line.fill.background()

# Title text block
close_box = slide12.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.0))
close_tf = close_box.text_frame
close_tf.word_wrap = True

p_c1 = close_tf.paragraphs[0]
p_c1.alignment = PP_ALIGN.CENTER
p_c1.text = "Join the Payments Revolution"
p_c1.font.name = "Trebuchet MS"
p_c1.font.size = Pt(48)
p_c1.font.bold = True
p_c1.font.color.rgb = TEXT_WHITE

p_c2 = close_tf.add_paragraph()
p_c2.alignment = PP_ALIGN.CENTER
p_c2.text = "Built on Stellar. Powered by Soroban."
p_c2.font.name = "Calibri"
p_c2.font.size = Pt(20)
p_c2.font.color.rgb = ACCENT_GOLD
p_c2.space_before = Pt(14)

p_c3 = close_tf.add_paragraph()
p_c3.alignment = PP_ALIGN.CENTER
p_c3.text = "Live Demo: stellar-escrowlevel-5.vercel.app  •  GitHub: github.com/shwetasharma44044-eng/StellarEscrowlevel-5"
p_c3.font.name = "Calibri"
p_c3.font.size = Pt(16)
p_c3.font.color.rgb = ACCENT_BLUE
p_c3.space_before = Pt(24)

p_c4 = close_tf.add_paragraph()
p_c4.alignment = PP_ALIGN.CENTER
p_c4.text = "Contact: contact@stellarescrow.com"
p_c4.font.name = "Calibri"
p_c4.font.size = Pt(14)
p_c4.font.color.rgb = TEXT_MUTED
p_c4.space_before = Pt(10)

# Save Presentation
prs.save("StellarEscrow_Pitch_Deck.pptx")
print("Pitch deck saved successfully as StellarEscrow_Pitch_Deck.pptx")
